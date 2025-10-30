# ppt_agent.py
from flask import Flask, request, jsonify
from pptx import Presentation
import os

app = Flask(__name__)

@app.route("/create_ppt", methods=["POST"])
def create_ppt():
    data = request.json
    summary = data.get("summary", "")
    filename = f"backend/storage/outputs/rfp_summary.pptx"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "RFP Summary"
    slide.placeholders[1].text = summary
    prs.save(filename)

    return jsonify({"ppt_path": filename})

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=8003)
